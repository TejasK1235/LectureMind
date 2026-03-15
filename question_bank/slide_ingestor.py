# question_bank/slide_ingestor.py

import os
import re
from typing import List, Dict, Optional


# ── Symbol normalisation map ───────────────────────────────────────────────
# Covers the most common math/CS symbols that survive PDF/PPTX extraction
# as Unicode but get garbled or dropped downstream in LLM context windows.
# Format: (compiled_regex_or_string, replacement)
_SYMBOL_MAP = [
    # Greek letters (common in CS/ML lectures)
    ("α", "alpha"), ("β", "beta"), ("γ", "gamma"), ("δ", "delta"),
    ("ε", "epsilon"), ("θ", "theta"), ("λ", "lambda"), ("μ", "mu"),
    ("π", "pi"), ("σ", "sigma"), ("τ", "tau"), ("φ", "phi"),
    ("ω", "omega"), ("Σ", "Sigma"), ("Δ", "Delta"), ("Ω", "Omega"),

    # Arrows
    ("→", "->"), ("←", "<-"), ("↔", "<->"), ("⇒", "=>"),
    ("⇐", "<="), ("⇔", "<=>"), ("↑", "up"), ("↓", "down"),

    # Math operators and relations
    ("≤", "<="), ("≥", ">="), ("≠", "!="), ("≈", "~="),
    ("∞", "infinity"), ("∑", "sum"), ("∏", "product"),
    ("∫", "integral"), ("∂", "partial"), ("∇", "gradient"),
    ("√", "sqrt"), ("∈", "in"), ("∉", "not in"),
    ("∩", "intersection"), ("∪", "union"), ("⊂", "subset of"),
    ("⊆", "subset or equal to"), ("∅", "empty set"),
    ("∀", "for all"), ("∃", "there exists"), ("¬", "not"),
    ("∧", "and"), ("∨", "or"), ("⊕", "XOR"),

    # Superscripts and subscripts (common in complexity notation)
    ("⁰","^0"),("¹","^1"),("²","^2"),("³","^3"),("⁴","^4"),
    ("⁵","^5"),("⁶","^6"),("⁷","^7"),("⁸","^8"),("⁹","^9"),
    ("₀","_0"),("₁","_1"),("₂","_2"),("₃","_3"),("₄","_4"),
    ("₅","_5"),("₆","_6"),("₇","_7"),("₈","_8"),("₉","_9"),

    # Fractions
    ("½", "1/2"), ("⅓", "1/3"), ("¼", "1/4"), ("¾", "3/4"),

    # Common CS notation
    ("∗", "*"), ("×", "x"), ("÷", "/"), ("±", "+/-"),
    ("′", "'"), ("″", "''"),

    # Bullets and list markers that survive as symbols
    ("•", "-"), ("·", "-"), ("◦", "-"), ("▪", "-"),
    ("▸", "->"), ("▶", "->"), ("►", "->"),

    # Quotes that get mangled
    ("\u201c", '"'), ("\u201d", '"'),
    ("\u2018", "'"), ("\u2019", "'"),

    # Dash variants
    ("\u2013", "-"), ("\u2014", "--"), ("\u2212", "-"),
]


def _normalise_symbols(text: str) -> str:
    """
    Replace Unicode math/CS symbols with readable ASCII equivalents.
    Also cleans up residual extraction artifacts like repeated spaces,
    stray newlines, and lone punctuation left behind by symbol removal.
    """
    for symbol, replacement in _SYMBOL_MAP:
        text = text.replace(symbol, replacement)

    # Collapse multiple spaces left behind by symbol removal
    text = re.sub(r" {2,}", " ", text)

    # Remove zero-width and non-printable characters that PDF extraction
    # sometimes leaves in (zero-width space, soft hyphen, etc.)
    text = re.sub(r"[\u200b\u200c\u200d\ufeff\u00ad]", "", text)

    # Clean up space before punctuation artifacts
    text = re.sub(r" ([,\.\;\:\!\?])", r"\1", text)

    return text.strip()


def ingest_slides(file_path: str) -> List[Dict]:
    """
    Extracts text from each slide/page of a PPT, PPTX, or PDF file.
    Returns a list of dicts:
    [
        {"slide_number": 1, "text": "extracted text from slide 1"},
        {"slide_number": 2, "text": "extracted text from slide 2"},
        ...
    ]
    Slides with no meaningful text are excluded.
    Symbol normalisation is applied to all extracted text so that
    math notation, Greek letters, and CS symbols are readable by the
    downstream LLM question generator.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Slide file not found: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()
    if ext in (".ppt", ".pptx"):
        return _ingest_pptx(file_path)
    elif ext == ".pdf":
        return _ingest_pdf(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}. Use PDF or PPTX.")


def _ingest_pptx(file_path: str) -> List[Dict]:
    """Extract text from each slide in a PPTX file."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation(file_path)
    slides = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = " ".join(
                    run.text.strip()
                    for run in para.runs
                    if run.text.strip()
                )
                if line:
                    texts.append(line)

        combined = _normalise_symbols(" ".join(texts).strip())

        if len(combined.split()) < 4:
            continue

        slides.append({
            "slide_number": i,
            "text": combined
        })

    return slides


def _ingest_pdf(file_path: str) -> List[Dict]:
    """Extract text from each page in a PDF file."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("PyMuPDF not installed. Run: pip install pymupdf")

    doc = fitz.open(file_path)
    slides = []

    for i, page in enumerate(doc, start=1):
        text = page.get_text().strip()
        text = " ".join(text.split())
        text = _normalise_symbols(text)

        if len(text.split()) < 4:
            continue

        slides.append({
            "slide_number": i,
            "text": text
        })

    doc.close()
    return slides








# import os
# from typing import List, Dict, Optional


# def ingest_slides(file_path: str) -> List[Dict]:
#     """
#     Extracts text from each slide/page of a PPT, PPTX, or PDF file.

#     Returns a list of dicts:
#     [
#         {"slide_number": 1, "text": "extracted text from slide 1"},
#         {"slide_number": 2, "text": "extracted text from slide 2"},
#         ...
#     ]

#     Slides with no meaningful text are excluded.
#     """

#     if not os.path.exists(file_path):
#         raise FileNotFoundError(f"Slide file not found: {file_path}")

#     ext = os.path.splitext(file_path)[1].lower()

#     if ext in (".ppt", ".pptx"):
#         return _ingest_pptx(file_path)
#     elif ext == ".pdf":
#         return _ingest_pdf(file_path)
#     else:
#         raise ValueError(f"Unsupported file format: {ext}. Use PDF or PPTX.")


# def _ingest_pptx(file_path: str) -> List[Dict]:
#     """Extract text from each slide in a PPTX file."""
#     try:
#         from pptx import Presentation
#     except ImportError:
#         raise ImportError("python-pptx not installed. Run: pip install python-pptx")

#     prs = Presentation(file_path)
#     slides = []

#     for i, slide in enumerate(prs.slides, start=1):
#         texts = []

#         for shape in slide.shapes:
#             # Only process shapes that have a text frame
#             if not shape.has_text_frame:
#                 continue

#             for para in shape.text_frame.paragraphs:
#                 line = " ".join(
#                     run.text.strip()
#                     for run in para.runs
#                     if run.text.strip()
#                 )
#                 if line:
#                     texts.append(line)

#         combined = " ".join(texts).strip()

#         # Skip slides with no meaningful text
#         if len(combined.split()) < 4:
#             continue

#         slides.append({
#             "slide_number": i,
#             "text": combined
#         })

#     return slides


# def _ingest_pdf(file_path: str) -> List[Dict]:
#     """Extract text from each page in a PDF file."""
#     try:
#         import fitz  # PyMuPDF
#     except ImportError:
#         raise ImportError("PyMuPDF not installed. Run: pip install pymupdf")

#     doc = fitz.open(file_path)
#     slides = []

#     for i, page in enumerate(doc, start=1):
#         text = page.get_text().strip()

#         # Clean up excessive whitespace and newlines
#         text = " ".join(text.split())

#         # Skip pages with no meaningful text
#         if len(text.split()) < 4:
#             continue

#         slides.append({
#             "slide_number": i,
#             "text": text
#         })

#     doc.close()
#     return slides